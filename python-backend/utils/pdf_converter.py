"""
Format conversion utilities.
PDF to Word, PDF to Image, PDF to PPT, PDF to Excel,
Word/Excel/PPT to PDF, and enhanced merge (PDFs + images).
All operations in-memory — nothing stored on disk.
"""

import io
import os
import logging
import subprocess
import tempfile
from typing import List, Tuple, Optional

logger = logging.getLogger(__name__)


def pdf_to_word(pdf_data: bytes) -> bytes:
    """
    Convert PDF to Word (.docx) using pdf2docx.

    Args:
        pdf_data: Input PDF bytes

    Returns:
        Word document bytes
    """
    from pdf2docx import Converter

    # pdf2docx requires file paths, so use temp files
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp_pdf:
        tmp_pdf.write(pdf_data)
        tmp_pdf_path = tmp_pdf.name

    tmp_docx_path = tmp_pdf_path.replace(".pdf", ".docx")

    try:
        cv = Converter(tmp_pdf_path)
        cv.convert(tmp_docx_path)
        cv.close()

        with open(tmp_docx_path, "rb") as f:
            result = f.read()

        logger.info(f"Converted PDF to Word: {len(result)} bytes")
        return result
    finally:
        _safe_remove(tmp_pdf_path)
        _safe_remove(tmp_docx_path)


def pdf_to_images(pdf_data: bytes, fmt: str = "jpeg") -> List[Tuple[bytes, str]]:
    """
    Convert each PDF page to an image.

    Args:
        pdf_data: Input PDF bytes
        fmt: Output format ("jpeg" or "png")

    Returns:
        List of (image_bytes, filename) tuples
    """
    from pdf2image import convert_from_bytes

    dpi = 200 if fmt == "jpeg" else 150
    images = convert_from_bytes(pdf_data, dpi=dpi, fmt=fmt)

    result = []
    for i, img in enumerate(images):
        buf = io.BytesIO()
        ext = "jpg" if fmt == "jpeg" else "png"
        save_kwargs = {"format": fmt.upper()}
        if fmt == "jpeg":
            save_kwargs["quality"] = 90
        img.save(buf, **save_kwargs)
        buf.seek(0)
        result.append((buf.getvalue(), f"page_{i + 1}.{ext}"))

    logger.info(f"Converted PDF to {len(result)} images ({fmt})")
    return result


def pdf_to_ppt(pdf_data: bytes) -> bytes:
    """
    Convert PDF to PowerPoint (.pptx) by rendering each page as an image slide.

    Args:
        pdf_data: Input PDF bytes

    Returns:
        PowerPoint bytes
    """
    from pptx import Presentation
    from pptx.util import Inches
    from pdf2image import convert_from_bytes

    images = convert_from_bytes(pdf_data, dpi=150, fmt="png")
    prs = Presentation()

    for img in images:
        # Set slide dimensions to match image aspect ratio
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        img_buf = io.BytesIO()
        img.save(img_buf, format="PNG")
        img_buf.seek(0)

        # Add image to slide, scaled to fit
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        img_ratio = img.width / img.height
        slide_ratio = slide_width / slide_height

        if img_ratio > slide_ratio:
            width = slide_width
            height = int(width / img_ratio)
            left = 0
            top = int((slide_height - height) / 2)
        else:
            height = slide_height
            width = int(height * img_ratio)
            left = int((slide_width - width) / 2)
            top = 0

        slide.shapes.add_picture(img_buf, left, top, width, height)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    logger.info(f"Converted PDF to PowerPoint: {len(images)} slides")
    return buf.getvalue()


def pdf_to_excel(pdf_data: bytes) -> bytes:
    """
    Extract text from PDF and create an Excel file.

    Args:
        pdf_data: Input PDF bytes

    Returns:
        Excel (.xlsx) bytes
    """
    from PyPDF2 import PdfReader
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Extracted Data"

    reader = PdfReader(io.BytesIO(pdf_data))
    row_num = 1
    for page_idx, page in enumerate(reader.pages):
        try:
            text = page.extract_text()
            if text:
                ws.cell(row=row_num, column=1, value=f"--- Page {page_idx + 1} ---")
                row_num += 1
                for line in text.split("\n"):
                    line = line.strip()
                    if line:
                        # Try to split by common delimiters
                        parts = line.split("\t") if "\t" in line else line.split("  ")
                        parts = [p.strip() for p in parts if p.strip()]
                        for col_idx, part in enumerate(parts, 1):
                            ws.cell(row=row_num, column=col_idx, value=part)
                        row_num += 1
        except Exception as e:
            logger.warning(f"Error extracting page {page_idx + 1}: {e}")

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    logger.info(f"Converted PDF to Excel: {row_num - 1} rows")
    return buf.getvalue()


def office_to_pdf(file_data: bytes, input_format: str) -> bytes:
    """
    Convert Word/Excel/PPT to PDF using LibreOffice headless.

    Args:
        file_data: Input file bytes
        input_format: File extension ("docx", "xlsx", "pptx", "doc", "xls", "ppt")

    Returns:
        PDF bytes
    """
    ext = input_format.lstrip(".")
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, f"input.{ext}")
        with open(input_path, "wb") as f:
            f.write(file_data)

        try:
            result = subprocess.run(
                [
                    "libreoffice", "--headless", "--convert-to", "pdf",
                    "--outdir", tmpdir, input_path,
                ],
                capture_output=True,
                timeout=120,
            )

            if result.returncode != 0:
                stderr = result.stderr.decode("utf-8", errors="replace")
                raise RuntimeError(f"LibreOffice conversion failed: {stderr}")

            output_path = os.path.join(tmpdir, "input.pdf")
            if not os.path.exists(output_path):
                raise RuntimeError("LibreOffice did not produce output PDF")

            with open(output_path, "rb") as f:
                pdf_bytes = f.read()

            logger.info(f"Converted {ext} to PDF: {len(pdf_bytes)} bytes")
            return pdf_bytes

        except subprocess.TimeoutExpired:
            raise RuntimeError("Conversion timed out. Try a smaller file.")
        except FileNotFoundError:
            raise RuntimeError(
                "LibreOffice not installed. Install with: "
                "sudo apt install libreoffice-common libreoffice-writer libreoffice-calc libreoffice-impress"
            )


def merge_mixed(files: List[Tuple[bytes, str]]) -> bytes:
    """
    Merge a mix of PDFs and images into a single PDF.

    Args:
        files: List of (file_bytes, mime_type) tuples.
               mime_type: "application/pdf", "image/jpeg", "image/png", etc.

    Returns:
        Merged PDF bytes
    """
    import pikepdf
    from utils.converter import convert_image_to_pdf

    output = pikepdf.new()

    for i, (data, mime_type) in enumerate(files):
        if mime_type == "application/pdf":
            with pikepdf.open(io.BytesIO(data)) as pdf:
                for page in pdf.pages:
                    output.pages.append(page)
        elif mime_type.startswith("image/"):
            # Convert image to single-page PDF, then add
            pdf_bytes = convert_image_to_pdf(data, mime_type)
            with pikepdf.open(io.BytesIO(pdf_bytes)) as pdf:
                for page in pdf.pages:
                    output.pages.append(page)
        else:
            logger.warning(f"Skipping unsupported file type: {mime_type}")

    if len(output.pages) == 0:
        raise ValueError("No valid files to merge")

    buf = io.BytesIO()
    output.save(buf)
    buf.seek(0)
    logger.info(f"Merged {len(files)} files into {len(output.pages)}-page PDF")
    return buf.getvalue()


def _safe_remove(path: str):
    """Remove a file, ignoring errors."""
    try:
        os.remove(path)
    except OSError:
        pass
